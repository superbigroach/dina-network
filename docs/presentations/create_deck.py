"""
Dina Network — Why Zero-Fee Blockchain Is Obvious
Presentation generator using python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)      # slate-950
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)      # slate-800
EMERALD = RGBColor(0x34, 0xD3, 0x99)      # emerald-400
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)         # slate-400
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)   # slate-300
RED = RGBColor(0xF8, 0x71, 0x71)          # red-400
AMBER = RGBColor(0xFB, 0xBF, 0x24)        # amber-400
BLUE = RGBColor(0x60, 0xA5, 0xFA)         # blue-400

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_bg(slide):
    """Fill slide background with dark color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text(slide, left, top, width, height, text, size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, size=22, color=WHITE, line_spacing=1.5, font_name="Segoe UI"):
    """Add multiple lines with consistent formatting."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, text_color, text_bold, text_size) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(text_size if text_size else size)
        p.font.color.rgb = text_color if text_color else color
        p.font.bold = text_bold if text_bold else False
        p.font.name = font_name
        p.space_after = Pt(size * 0.6)
    return txBox

def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_text(slide, 2, 1.5, 12, 1, "DINA NETWORK", 20, EMERALD, True)
add_text(slide, 2, 2.3, 12, 2, "Why Every Blockchain Charges\nFees — And Why We Don't", 52, WHITE, True)
add_text(slide, 2, 5.0, 12, 1, "The $0.00 Transaction", 28, EMERALD, False)
add_text(slide, 2, 6.5, 12, 1, "100ms finality  ·  Zero fees  ·  USDC native  ·  21 validators", 18, GRAY)
add_text(slide, 2, 7.5, 12, 1, "dina-wallet.web.app", 16, GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: The Question
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 1, 12, 1, "THE QUESTION", 18, EMERALD, True)
add_text(slide, 2, 2, 12, 2, "Why does it cost money\nto send money?", 52, WHITE, True)
add_multiline(slide, 2, 4.8, 12, 4, [
    ("Ethereum: $1 – $50 per transaction", RED, False, 28),
    ("Solana: $0.001 – but you need SOL tokens first", AMBER, False, 28),
    ("Base: $0.01 – but you need ETH tokens first", AMBER, False, 28),
    ("", WHITE, False, 14),
    ("Dina: $0.00 — just USDC, nothing else needed", EMERALD, True, 32),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: What fees actually pay for
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 1, 12, 1, "WHAT GAS FEES ACTUALLY PAY FOR", 18, EMERALD, True)
add_text(slide, 2, 2, 12, 1.5, "Rent. For other people's computers.", 48, WHITE, True)

add_card(slide, 2, 4, 5.5, 3.5)
add_text(slide, 2.4, 4.2, 5, 0.5, "EVERY OTHER BLOCKCHAIN", 14, GRAY, True)
add_multiline(slide, 2.4, 4.8, 5, 2.8, [
    ("Validators = strangers", WHITE, False, 22),
    ("Strangers won't work for free", WHITE, False, 22),
    ("So you charge users gas fees", WHITE, False, 22),
    ("Fees go to the strangers", WHITE, False, 22),
    ("→ Users pay $1-50 per transaction", RED, True, 22),
])

add_card(slide, 8.5, 4, 5.5, 3.5)
add_text(slide, 8.9, 4.2, 5, 0.5, "DINA", 14, EMERALD, True)
add_multiline(slide, 8.9, 4.8, 5, 2.8, [
    ("Validators = our servers", WHITE, False, 22),
    ("We pay the cloud bill", WHITE, False, 22),
    ("Users pay nothing", WHITE, False, 22),
    ("Same cryptographic security", WHITE, False, 22),
    ("→ Users pay $0.00 per transaction", EMERALD, True, 22),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: The Rube Goldberg Machine
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "WHAT OTHER CHAINS BUILT TO PAY STRANGERS", 18, EMERALD, True)

items = [
    ("Native Token (ETH, SOL, BNB)", "Medium of payment to validators"),
    ("Gas Fee Market", "Dynamic pricing so blocks don't get spammed"),
    ("Staking Mechanism", "Force validators to lock capital as collateral"),
    ("Slashing Rules", "Punish misbehaving validators"),
    ("Inflation Schedule", "Print new tokens when fees aren't enough"),
    ("Token Burns (EIP-1559)", "Burn tokens to control supply inflation"),
    ("MEV Extraction", "Validators reorder txs for profit (side effect)"),
    ("MEV Protection (Flashbots)", "Protect users from validator exploitation"),
    ("Governance Systems", "Let token holders vote on fee parameters"),
    ("Liquid Staking (Lido)", "Let people stake without running a validator"),
]

for i, (title, desc) in enumerate(items):
    y = 1.8 + i * 0.65
    add_text(slide, 2, y, 0.5, 0.5, str(i+1), 16, EMERALD, True)
    add_text(slide, 2.6, y, 5, 0.5, title, 18, WHITE, True)
    add_text(slide, 8, y, 6, 0.5, desc, 16, GRAY)

add_text(slide, 2, 8.3, 12, 0.5, "All of this exists because validators are anonymous strangers who need to be paid.", 18, AMBER, True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: Dina's Solution
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 1, 12, 1, "DINA'S SOLUTION", 18, EMERALD, True)
add_text(slide, 2, 2.2, 12, 1.5, "Just run the servers yourself.", 52, WHITE, True)

add_card(slide, 2, 4.5, 12, 1.2)
add_text(slide, 2.5, 4.55, 11, 1.2, "21 validators  ×  Google Cloud e2-medium  ×  $24/month each  =  $504/month total", 28, EMERALD, True, PP_ALIGN.CENTER)

add_text(slide, 2, 6.2, 12, 0.5, "That's the entire network. $504/month.", 24, WHITE, True)
add_text(slide, 2, 7.0, 12, 2, "No token. No staking. No inflation. No fee market. No governance.\nNo MEV. No liquid staking. No slashing. No gas fees.\nJust 21 servers verifying transactions.", 22, GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: "But what do validators actually DO?"
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "WHAT VALIDATORS ACTUALLY DO", 18, EMERALD, True)
add_text(slide, 2, 1.6, 12, 1, "The same thing on every blockchain. Including ours.", 28, GRAY)

steps = [
    ("1", "Receive transaction from user"),
    ("2", "Verify Ed25519 signature (is this really from the sender?)"),
    ("3", "Check nonce (prevent replay attacks)"),
    ("4", "Check balance (does sender have enough USDC?)"),
    ("5", "Execute transaction (move money A → B)"),
    ("6", "Run gas meter on smart contracts (prevent infinite loops)"),
    ("7", "Update state (new balances)"),
    ("8", "Package into a block, sign it"),
    ("9", "Send to other validators, reach consensus"),
    ("10", "Block confirmed — 100ms total"),
]

for i, (num, desc) in enumerate(steps):
    y = 2.8 + i * 0.58
    add_text(slide, 2, y, 0.6, 0.5, num, 18, EMERALD, True)
    add_text(slide, 2.8, y, 11, 0.5, desc, 20, WHITE if i < 9 else EMERALD, i == 9)

add_text(slide, 2, 8.6, 12, 0.5, "Identical process on Ethereum, Solana, BSC, and Dina. The only difference: who pays the electricity.", 16, GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: Cost Comparison
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "WHAT IT COSTS TO RUN A BLOCKCHAIN", 18, EMERALD, True)

# Table header
add_card(slide, 1.5, 1.8, 13, 0.7)
headers = [("Chain", 1.7, 2.5), ("RAM", 4.2, 1.5), ("Cost/Validator", 5.8, 2.5), ("# Validators", 8.5, 2), ("User Fee", 10.7, 2.5)]
for text, x, w in headers:
    add_text(slide, x, 1.85, w, 0.5, text, 16, EMERALD, True)

rows = [
    ("Solana", "512 GB", "$6-12K/mo", "1,800", "$0.001 + need SOL", RED),
    ("Ethereum", "16-32 GB", "$80-150/mo", "1,050,000", "$1-50 + need ETH", RED),
    ("BSC", "64-128 GB", "$300-600/mo", "21", "$0.05 + need BNB", AMBER),
    ("Avalanche", "16-32 GB", "$100-300/mo", "1,200", "$0.01 + need AVAX", AMBER),
    ("Polygon", "32-64 GB", "$200-500/mo", "100", "$0.01 + need POL", AMBER),
    ("Dina", "4 GB", "$24/mo", "21", "$0.00", EMERALD),
]

for i, (chain, ram, cost, vals, fee, fee_color) in enumerate(rows):
    y = 2.7 + i * 0.75
    if i == len(rows) - 1:
        add_card(slide, 1.5, y - 0.1, 13, 0.65, RGBColor(0x06, 0x4E, 0x3B))
    add_text(slide, 1.7, y, 2.5, 0.5, chain, 20, WHITE if i < len(rows)-1 else EMERALD, i == len(rows)-1)
    add_text(slide, 4.2, y, 1.5, 0.5, ram, 18, GRAY if i < len(rows)-1 else EMERALD)
    add_text(slide, 5.8, y, 2.5, 0.5, cost, 18, GRAY if i < len(rows)-1 else EMERALD)
    add_text(slide, 8.5, y, 2, 0.5, vals, 18, GRAY if i < len(rows)-1 else EMERALD)
    add_text(slide, 10.7, y, 2.5, 0.5, fee, 18, fee_color, i == len(rows)-1)

add_text(slide, 2, 7.5, 12, 1, "Solana validators pay $5,000-10,000/month just in voting fees.\nEthereum validators must lock $80,000+ in ETH as collateral.\nDina validators cost less than a Netflix subscription.", 20, GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: Why so expensive
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "WHY DO OTHER CHAINS NEED SUCH EXPENSIVE HARDWARE?", 18, EMERALD, True)

add_card(slide, 1.5, 1.8, 6, 3.2)
add_text(slide, 1.8, 1.9, 5.5, 0.5, "SOLANA — $512 GB RAM", 18, RED, True)
add_multiline(slide, 1.8, 2.5, 5.5, 2.5, [
    ("Keeps entire accounts database", WHITE, False, 18),
    ("(~400 GB) in memory", WHITE, False, 18),
    ("", WHITE, False, 10),
    ("Every validator replays every", WHITE, False, 18),
    ("transaction at 4,000 TPS", WHITE, False, 18),
    ("", WHITE, False, 10),
    ("Needs 1 Gbps bandwidth minimum", WHITE, False, 18),
])

add_card(slide, 8.5, 1.8, 6, 3.2)
add_text(slide, 8.8, 1.9, 5.5, 0.5, "ETHEREUM — 1,050,000 validators", 18, RED, True)
add_multiline(slide, 8.8, 2.5, 5.5, 2.5, [
    ("$80,000+ locked per validator", WHITE, False, 18),
    ("= ~$84 BILLION total staked", WHITE, False, 18),
    ("", WHITE, False, 10),
    ("Must process attestations from", WHITE, False, 18),
    ("1 million validators per epoch", WHITE, False, 18),
    ("", WHITE, False, 10),
    ("2-4 TB SSD for chain state", WHITE, False, 18),
])

add_card(slide, 1.5, 5.3, 13, 2.8)
add_text(slide, 1.8, 5.4, 12.5, 0.5, "DINA — 4 GB RAM, $24/month", 20, EMERALD, True)
add_multiline(slide, 1.8, 6.0, 12, 2, [
    ("Fresh chain = minimal state to store", WHITE, False, 20),
    ("21 known validators = trivial consensus overhead", WHITE, False, 20),
    ("No voting fees, no staking collateral", WHITE, False, 20),
    ("Same Ed25519 cryptographic verification as everyone else", WHITE, False, 20),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: Gas metering vs gas fees
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "GAS METERING WITHOUT GAS FEES", 18, EMERALD, True)
add_text(slide, 2, 1.7, 12, 1, "We track gas for safety. We charge $0 for it.", 32, WHITE, True)

add_card(slide, 1.5, 3.2, 6, 4)
add_text(slide, 1.8, 3.3, 5.5, 0.5, "OTHER CHAINS USE GAS FOR TWO THINGS", 14, RED, True)
add_multiline(slide, 1.8, 3.9, 5.5, 3, [
    ("1. Safety", WHITE, True, 22),
    ("   Prevent infinite loops in", GRAY, False, 18),
    ("   smart contracts", GRAY, False, 18),
    ("", WHITE, False, 10),
    ("2. Billing", WHITE, True, 22),
    ("   Charge users for computation", GRAY, False, 18),
    ("   Route payment to validators", GRAY, False, 18),
])

add_card(slide, 8.5, 3.2, 6, 4)
add_text(slide, 8.8, 3.3, 5.5, 0.5, "DINA SEPARATES THEM", 14, EMERALD, True)
add_multiline(slide, 8.8, 3.9, 5.5, 3, [
    ("Gas Metering: ON", EMERALD, True, 22),
    ("   Every operation is counted", GRAY, False, 18),
    ("   Infinite loops get killed", GRAY, False, 18),
    ("", WHITE, False, 10),
    ("Gas Price: $0.00", EMERALD, True, 22),
    ("   Meter runs but costs nothing", GRAY, False, 18),
    ("   Safety without the bill", GRAY, False, 18),
])

add_text(slide, 2, 7.6, 12, 0.7, "Like a car with a rev limiter but free fuel.\nThe engine has safety limits. The gas station doesn't exist.", 20, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: Spam prevention
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, '"WITHOUT FEES, WON\'T PEOPLE SPAM THE NETWORK?"', 18, EMERALD, True)
add_text(slide, 2, 1.7, 12, 1, "No. We use balance-based rate limiting.", 32, WHITE, True)

add_card(slide, 1.5, 3, 13, 1)
add_text(slide, 1.8, 3, 12.5, 0.5, "$1 minimum balance to transact (Sybil attack = hold $1M USDC = expensive)", 20, WHITE, False, PP_ALIGN.CENTER)

rows = [
    ("$1+", "100 burst", "1 tx/sec"),
    ("$100+", "500 burst", "5 tx/sec"),
    ("$1,000+", "2,000 burst", "20 tx/sec"),
    ("$10,000+", "10,000 burst", "100 tx/sec"),
    ("$100,000+", "50,000 burst", "500 tx/sec"),
]

add_text(slide, 2, 4.3, 3, 0.5, "Balance", 16, EMERALD, True)
add_text(slide, 5.5, 4.3, 3, 0.5, "Burst Capacity", 16, EMERALD, True)
add_text(slide, 9, 4.3, 3, 0.5, "Sustained Rate", 16, EMERALD, True)

for i, (bal, burst, rate) in enumerate(rows):
    y = 4.9 + i * 0.6
    add_text(slide, 2, y, 3, 0.5, bal, 22, WHITE, True)
    add_text(slide, 5.5, y, 3, 0.5, burst, 22, GRAY)
    add_text(slide, 9, y, 3, 0.5, rate, 22, GRAY)

add_text(slide, 2, 8, 12, 0.7, "Token bucket algorithm. More money = more capacity. No fees needed.", 20, GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: 9-wallet system
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "EVERY USER GETS 9 WALLETS", 18, EMERALD, True)
add_text(slide, 2, 1.6, 12, 1, "Three tiers. Independent keys. All earning yield.", 28, GRAY)

# SCA tier
add_card(slide, 1.5, 2.7, 4, 4.5)
add_text(slide, 1.8, 2.8, 3.5, 0.5, "SMART ACCOUNTS (x3)", 14, EMERALD, True)
add_text(slide, 1.8, 3.3, 3.5, 0.5, "4.5% APY", 24, EMERALD, True)
add_multiline(slide, 1.8, 4, 3.5, 3, [
    ("Main Wallet", WHITE, True, 20),
    ("Savings", WHITE, True, 20),
    ("Backup", WHITE, True, 20),
    ("", WHITE, False, 10),
    ("You control the keys", GRAY, False, 16),
    ("Each has its own password", GRAY, False, 16),
    ("Can freeze other wallets", GRAY, False, 16),
])

# Agent tier
add_card(slide, 6, 2.7, 4, 4.5)
add_text(slide, 6.3, 2.8, 3.5, 0.5, "AGENT WALLETS (x3)", 14, AMBER, True)
add_text(slide, 6.3, 3.3, 3.5, 0.5, "3.5% APY", 24, AMBER, True)
add_multiline(slide, 6.3, 4, 3.5, 3, [
    ("Shopping Bot", WHITE, True, 20),
    ("Bill Payments", WHITE, True, 20),
    ("Custom Agent", WHITE, True, 20),
    ("", WHITE, False, 10),
    ("AI signs within daily limits", GRAY, False, 16),
    ("Revocable from any SCA", GRAY, False, 16),
    ("$200-$2,000/day caps", GRAY, False, 16),
])

# Parallel tier
add_card(slide, 10.5, 2.7, 4, 4.5)
add_text(slide, 10.8, 2.8, 3.5, 0.5, "PARALLEL WALLETS (x3)", 14, BLUE, True)
add_text(slide, 10.8, 3.3, 3.5, 0.5, "3.0% APY", 24, BLUE, True)
add_multiline(slide, 10.8, 4, 3.5, 3, [
    ("Business Payments", WHITE, True, 20),
    ("Streaming Payments", WHITE, True, 20),
    ("Developer API", WHITE, True, 20),
    ("", WHITE, False, 10),
    ("Independent nonces", GRAY, False, 16),
    ("Simultaneous tx streams", GRAY, False, 16),
    ("1,000 payments in 100ms", GRAY, False, 16),
])

add_text(slide, 2, 7.8, 12, 0.7, "All USDC. All earning yield. Users keep 100% of their yield.", 22, EMERALD, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: Why nobody did this
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "WHY HASN'T ANYONE DONE THIS BEFORE?", 18, EMERALD, True)

reasons = [
    ("1", "Tokens make founders rich", "Create a token, keep 20-40% of supply. Token goes up = billionaire.\nSolana Foundation holds ~300M SOL. Why would they eliminate fees?", RED),
    ("2", "VCs fund tokens, not infrastructure", "Invest $10M → get tokens → token launches → 100x return.\n\"We just run servers\" doesn't give VCs a token to flip.", AMBER),
    ("3", "\"Decentralization\" is a shield", "Question the fees? You're questioning decentralization.\nNobody examines it because everyone's making money.", AMBER),
    ("4", "Most people don't understand it", "99% of crypto users don't know a validator is just a server.\n\"Decentralization requires fees\" is accepted like gravity.", GRAY),
    ("5", "The people who know are already rich", "Engineers at Ethereum, Solana, Coinbase know all of this.\nZero incentive to build the simpler version.", GRAY),
]

for i, (num, title, desc, color) in enumerate(reasons):
    y = 1.8 + i * 1.35
    add_text(slide, 2, y, 0.5, 0.5, num, 24, color, True)
    add_text(slide, 2.7, y, 11, 0.5, title, 24, WHITE, True)
    add_text(slide, 2.7, y + 0.5, 11, 0.8, desc, 16, GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13: The comparison
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, '"BUT ISN\'T THAT CENTRALIZED?"', 18, EMERALD, True)

props = [
    ("Open-source code", True, True),
    ("Public ledger (anyone can read)", True, True),
    ("Cryptographic verification", True, True),
    ("Users sign own transactions", True, True),
    ("Anyone can audit the chain", True, True),
    ("Anyone can run a full node", True, True),
    ("Anyone can fork the code", True, True),
    ("Validators are anonymous strangers", True, False),
    ("Users pay fees", True, False),
]

add_text(slide, 5.5, 1.7, 3, 0.5, "Ethereum", 18, GRAY, True, PP_ALIGN.CENTER)
add_text(slide, 9, 1.7, 3, 0.5, "Dina", 18, EMERALD, True, PP_ALIGN.CENTER)

for i, (prop, eth, dina) in enumerate(props):
    y = 2.3 + i * 0.62
    add_text(slide, 2, y, 3.3, 0.5, prop, 18, WHITE)
    add_text(slide, 6.2, y, 1, 0.5, "Yes" if eth else "No", 18, EMERALD if eth else RED, False, PP_ALIGN.CENTER)
    color = EMERALD if dina else AMBER
    add_text(slide, 9.7, y, 1, 0.5, "Yes" if dina else "No", 18, color, dina != eth, PP_ALIGN.CENTER)

add_text(slide, 2, 8.2, 12, 0.5, "Same transparency. Same crypto. Same auditability. Minus the fees.", 20, EMERALD, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14: The real numbers
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "THE BUSINESS MODEL", 18, EMERALD, True)
add_text(slide, 2, 1.6, 12, 1, "It costs almost nothing. Revenue comes from everywhere else.", 28, GRAY)

add_card(slide, 1.5, 2.8, 6, 2.5)
add_text(slide, 1.8, 2.9, 5.5, 0.5, "COSTS", 14, RED, True)
add_multiline(slide, 1.8, 3.4, 5.5, 2, [
    ("Testnet (now): $75/month", WHITE, False, 22),
    ("Mainnet (21 validators): $504/month", WHITE, False, 22),
    ("At scale (100 TPS): ~$6,600/month", WHITE, False, 22),
    ("At Visa-level (1K TPS): ~$17K/month", WHITE, False, 22),
])

add_card(slide, 8.5, 2.8, 6, 2.5)
add_text(slide, 8.8, 2.9, 5.5, 0.5, "REVENUE", 14, EMERALD, True)
add_multiline(slide, 8.8, 3.4, 5.5, 2, [
    ("$1 wallet activation per user", WHITE, False, 22),
    ("Enterprise tiers: $99-$999/mo", WHITE, False, 22),
    ("Treasury yield: 4.5% on reserves", WHITE, False, 22),
    ("500K users = $500K (covers everything)", WHITE, False, 22),
])

add_card(slide, 1.5, 5.8, 13, 2.2)
add_text(slide, 1.8, 5.9, 12.5, 0.5, "USER ECONOMICS", 14, EMERALD, True)
add_multiline(slide, 1.8, 6.4, 12, 1.5, [
    ("Transactions: $0.00 (zero, always)", EMERALD, True, 24),
    ("Yield: 4.5% APY on USDC — user keeps 100%", EMERALD, True, 24),
    ("Speed: 100ms finality (faster than tapping your credit card)", WHITE, False, 20),
    ("No gas token needed. No ETH. No SOL. Just USDC.", WHITE, False, 20),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15: Final slide
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 2, 1.5, 12, 2, "Every other blockchain charges fees\nbecause they built a system that\nrequires paying strangers.", 40, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, 2, 4.5, 12, 1.5, "We don't have strangers.\nWe have a cloud bill.", 48, WHITE, True, PP_ALIGN.CENTER)

add_text(slide, 2, 7, 12, 0.8, "$0.00 per transaction  ·  100ms finality  ·  4.5% APY on your USDC", 24, EMERALD, True, PP_ALIGN.CENTER)
add_text(slide, 2, 8, 12, 0.5, "dina-wallet.web.app", 20, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_path = "C:/dina_network/docs/presentations/Dina_Zero_Fee_Blockchain.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
