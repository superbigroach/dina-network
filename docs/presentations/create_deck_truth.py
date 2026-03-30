"""
The Crypto Fee Scam — How Blockchain Became a Toll Booth
A presentation about the uncomfortable truths of crypto economics
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
EMERALD = RGBColor(0x34, 0xD3, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
RED = RGBColor(0xF8, 0x71, 0x71)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text(slide, left, top, width, height, text, size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, size=22, color=WHITE, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, text_color, text_bold, text_size) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(text_size if text_size else size)
        p.font.color.rgb = text_color if text_color else color
        p.font.bold = text_bold if text_bold else False
        p.font.name = font_name
        p.space_after = Pt(size * 0.5)
    return txBox

def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 2, 12, 2, "The Blockchain\nFee Machine", 56, WHITE, True)
add_text(slide, 2, 4.8, 12, 1, "How an industry convinced the world that\nmoving numbers in a database should cost $50", 28, GRAY)
add_text(slide, 2, 7, 12, 0.5, "And why they'll never fix it — because they're getting rich from it.", 20, RED)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: What a transaction actually is
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "WHAT A BLOCKCHAIN TRANSACTION ACTUALLY IS", 18, EMERALD, True)

add_text(slide, 2, 2, 12, 1, "This is what happens when you send $10:", 28, WHITE)

add_card(slide, 2, 3.2, 12, 4)
add_multiline(slide, 2.5, 3.4, 11, 3.5, [
    ("1. Check signature (is this really from you?)                    ~0.001ms", GRAY, False, 20),
    ("2. Check nonce (has this tx been seen before?)                   ~0.0001ms", GRAY, False, 20),
    ("3. Check balance (do you have $10?)                                   ~0.001ms", GRAY, False, 20),
    ("4. Subtract $10 from your balance                                        ~0.0001ms", GRAY, False, 20),
    ("5. Add $10 to recipient's balance                                          ~0.0001ms", GRAY, False, 20),
    ("6. Write to disk                                                                        ~0.01ms", GRAY, False, 20),
    ("", WHITE, False, 10),
    ("Total computation: ~0.01 milliseconds", WHITE, True, 24),
    ("Total electricity: ~$0.000001", WHITE, True, 24),
])

add_text(slide, 2, 7.8, 12, 0.8, "Ethereum charges $1-50 for this.\nThat's a 1,000,000x to 50,000,000x markup on the actual cost.", 24, RED, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: Where the money goes
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "SO WHERE DOES YOUR $50 GAS FEE GO?", 18, EMERALD, True)
add_text(slide, 2, 1.7, 12, 1, "Not to computation. Not to storage. Not to security.", 28, GRAY)

add_text(slide, 2, 3, 12, 1, "It goes to strangers running servers.", 36, WHITE, True)

add_multiline(slide, 2, 4.2, 12, 4, [
    ("The entire gas fee model exists for one reason:", WHITE, False, 24),
    ("", WHITE, False, 10),
    ("Ethereum has ~1,050,000 validators run by anonymous people.", LIGHT_GRAY, False, 22),
    ("Anonymous people don't volunteer their computers for free.", LIGHT_GRAY, False, 22),
    ("So you — the user — pay them every time you move your own money.", LIGHT_GRAY, False, 22),
    ("", WHITE, False, 10),
    ("You are paying rent on someone else's computer", RED, True, 28),
    ("to update a number in a database.", RED, True, 28),
    ("", WHITE, False, 10),
    ("Every single time.", RED, False, 22),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: The Token Scam
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "THE TOKEN GAME", 18, RED, True)
add_text(slide, 2, 1.6, 12, 1, "Why every blockchain has its own token", 36, WHITE, True)

add_text(slide, 2, 3, 12, 0.6, "The stated reason:", 20, GRAY)
add_text(slide, 2, 3.5, 12, 0.8, '"We need a native token to align validator incentives\nand secure the network through staking."', 24, LIGHT_GRAY, False)

add_text(slide, 2, 4.8, 12, 0.6, "The actual reason:", 20, AMBER)

add_card(slide, 2, 5.4, 12, 3)
add_multiline(slide, 2.5, 5.5, 11, 2.8, [
    ("Step 1: Create a token. Keep 20-40% for founders + insiders.", RED, True, 22),
    ("Step 2: Sell tokens to VCs at a discount.", RED, False, 22),
    ("Step 3: Launch the chain. Require the token for gas fees.", RED, False, 22),
    ("Step 4: Every user must buy the token to do anything.", RED, False, 22),
    ("Step 5: Demand pushes price up.", RED, False, 22),
    ("Step 6: Founders and VCs sell into that demand.", RED, True, 22),
    ("", WHITE, False, 10),
    ("Gas fees aren't a technical necessity. They're a monetization strategy.", AMBER, True, 22),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: Founder Wealth
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "HOW MUCH FOUNDERS MAKE FROM THE TOKEN MODEL", 18, RED, True)

add_card(slide, 1.5, 1.8, 13, 1.5)
add_text(slide, 2, 1.9, 4, 0.5, "Chain", 16, GRAY, True)
add_text(slide, 5.5, 1.9, 3, 0.5, "Founder/Foundation Holdings", 16, GRAY, True)
add_text(slide, 9.5, 1.9, 4, 0.5, "Estimated Value", 16, GRAY, True)
add_text(slide, 2, 2.5, 4, 0.5, "Ethereum", 22, WHITE, True)
add_text(slide, 5.5, 2.5, 3.5, 0.5, "Foundation: ~300K ETH", 18, GRAY)
add_text(slide, 9.5, 2.5, 4, 0.5, "$600M - $1B+", 22, RED, True)

rows = [
    ("Solana", "Foundation + Labs: ~300M SOL", "$15B - $90B", "At peak, more than Ford Motor Company"),
    ("BNB", "Binance: ~80M BNB", "$30B - $50B", "CZ was richest person in crypto"),
    ("Cardano", "Foundation + IOHK: ~6B ADA", "$2B - $10B", "Charles Hoskinson: yacht, ranch"),
    ("Ripple", "Company: ~40B XRP", "$20B - $100B", "Founders among richest in crypto"),
    ("Avalanche", "Foundation: ~270M AVAX", "$3B - $15B", "Token sale raised $60M in hours"),
]

for i, (chain, holdings, value, note) in enumerate(rows):
    y = 3.4 + i * 0.85
    add_text(slide, 2, y, 3, 0.5, chain, 22, WHITE, True)
    add_text(slide, 5.5, y, 3.5, 0.5, holdings, 16, GRAY)
    add_text(slide, 9.5, y, 2, 0.5, value, 20, RED, True)
    add_text(slide, 11.7, y, 3, 0.5, note, 12, GRAY)

add_text(slide, 2, 7.8, 12, 0.8, "These people have zero incentive to eliminate gas fees.\nFees create demand for the token. Demand makes their holdings worth billions.", 20, AMBER, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: The VC Game
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "THE VC PLAYBOOK", 18, RED, True)
add_text(slide, 2, 1.6, 12, 1, "Why crypto VCs only fund token projects", 32, WHITE, True)

add_card(slide, 1.5, 2.8, 6.2, 3.5)
add_text(slide, 1.8, 2.9, 5.8, 0.5, "TOKEN PROJECT (what VCs fund)", 14, RED, True)
add_multiline(slide, 1.8, 3.5, 5.8, 2.8, [
    ("VC invests $10M at token discount", WHITE, False, 20),
    ("Gets tokens at $0.10 each", WHITE, False, 20),
    ("Token launches at $1.00", WHITE, False, 20),
    ("VC sells for $100M", AMBER, True, 20),
    ("Return: 10x in 6-12 months", RED, True, 24),
    ("", WHITE, False, 10),
    ("VCs don't care if the tech works.", GRAY, False, 16),
    ("They care if the token pumps.", GRAY, False, 16),
])

add_card(slide, 8.3, 2.8, 6.2, 3.5)
add_text(slide, 8.6, 2.9, 5.8, 0.5, "INFRASTRUCTURE COMPANY (what Dina is)", 14, EMERALD, True)
add_multiline(slide, 8.6, 3.5, 5.8, 2.8, [
    ("VC invests $5M for equity", WHITE, False, 20),
    ("Company grows revenue over years", WHITE, False, 20),
    ("Exit at $50-100M in 5-7 years", WHITE, False, 20),
    ("Return: 10-20x over 5+ years", EMERALD, True, 20),
    ("Return: based on real revenue", EMERALD, True, 24),
    ("", WHITE, False, 10),
    ("Same returns, but requires patience", GRAY, False, 16),
    ("and a product that actually works.", GRAY, False, 16),
])

add_text(slide, 2, 6.8, 12, 1.5, "The entire crypto VC ecosystem is optimized for token speculation,\nnot for building products people use.\n\nThat's why every chain has a token. Not because they need one.", 22, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: The Decentralization Lie
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, 'THE "DECENTRALIZATION" SMOKESCREEN', 18, RED, True)
add_text(slide, 2, 1.6, 12, 1, "They say it's about decentralization.\nLook at who actually controls these networks.", 28, GRAY)

chains = [
    ("Solana", "Top 19 validators control 33% of stake (halt threshold).\nSolana Labs + Foundation effectively control upgrades.\nChain has halted 10+ times — centralized restart each time.", RED),
    ("BSC", "21 validators. All selected by BNB staking.\nBinance controls enough BNB to select every validator.\nIt's Binance's private database with extra steps.", RED),
    ("Base", "ONE sequencer. Run by Coinbase.\nLiterally a single company's server.\nBut they still charge gas fees in ETH.", AMBER),
    ("Polygon", "~100 validators. Top 10 control majority of stake.\nPolygon Labs controls the upgrade process.\nMultisig with known insiders can override the chain.", AMBER),
]

for i, (chain, desc, color) in enumerate(chains):
    y = 3.2 + i * 1.45
    add_card(slide, 1.5, y - 0.1, 13, 1.3)
    add_text(slide, 1.8, y, 2, 0.5, chain, 22, color, True)
    add_text(slide, 4, y, 10.2, 1.1, desc, 16, GRAY)

add_text(slide, 2, 8.3, 12, 0.5, "\"Decentralized\" is a marketing word, not a technical description of these networks.", 20, AMBER, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: Validators aren't altruists
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "VALIDATORS AREN'T VOLUNTEERS. THEY'RE BUSINESSES.", 18, RED, True)

add_text(slide, 2, 1.8, 12, 1, "They run validators to make money. That's the only reason.", 28, WHITE, True)

add_card(slide, 1.5, 3, 13, 2)
add_text(slide, 1.8, 3.1, 12.5, 0.5, "SOLANA VALIDATOR ECONOMICS", 14, RED, True)
add_multiline(slide, 1.8, 3.6, 12, 1.5, [
    ("Hardware: $1,000-1,800/month  |  Voting fees: $5,000-10,000/month  |  Total: ~$7-12K/month", WHITE, False, 18),
    ("Revenue: inflation rewards + MEV extraction + transaction fees", WHITE, False, 18),
    ("Many validators are unprofitable. They run at a loss hoping SOL price goes up.", AMBER, False, 18),
    ("This is not decentralized security. This is speculation on a token.", AMBER, True, 18),
])

add_card(slide, 1.5, 5.3, 13, 2)
add_text(slide, 1.8, 5.4, 12.5, 0.5, "ETHEREUM VALIDATOR ECONOMICS", 14, RED, True)
add_multiline(slide, 1.8, 5.9, 12, 1.5, [
    ("Must lock 32 ETH (~$80,000+) per validator as collateral", WHITE, False, 18),
    ("Earn ~4-5% APY in ETH rewards (paid by inflating the supply)", WHITE, False, 18),
    ("Lido controls ~28% of all staked ETH. So much for decentralization.", AMBER, False, 18),
    ("$80B in capital locked up doing nothing productive. Just sitting there.", AMBER, True, 18),
])

add_text(slide, 2, 7.8, 12, 1, "Validators aren't securing the network out of idealism.\nThey're running a business that profits from your gas fees and token inflation.", 22, RED, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: What you actually need a token for
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "DO YOU ACTUALLY NEED A NATIVE TOKEN?", 18, EMERALD, True)
add_text(slide, 2, 1.8, 12, 1, "No.", 72, WHITE, True)

add_text(slide, 2, 3.5, 12, 0.6, "What they tell you tokens are for:", 22, GRAY, True)
add_multiline(slide, 2, 4.2, 12, 3, [
    ('"Pay validators"          → Pay them from revenue, like every other company', LIGHT_GRAY, False, 20),
    ('"Secure the network"     → 21 known validators with SLAs. Done.', LIGHT_GRAY, False, 20),
    ('"Governance"                  → A company makes decisions. That\'s called management.', LIGHT_GRAY, False, 20),
    ('"Align incentives"           → Employees have incentives. It\'s called a salary.', LIGHT_GRAY, False, 20),
    ('"Decentralization"          → BSC has 21 validators. Base has 1. They have tokens anyway.', LIGHT_GRAY, False, 20),
])

add_card(slide, 2, 7.2, 12, 1.3)
add_multiline(slide, 2.5, 7.3, 11, 1, [
    ("The real reason tokens exist:", AMBER, True, 22),
    ("To make founders and VCs rich through artificial demand created by mandatory gas fees.", RED, True, 22),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: The fee complexity
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "THE ABSURD COMPLEXITY OF BLOCKCHAIN FEES", 18, RED, True)
add_text(slide, 2, 1.6, 12, 1, "Things a normal person has to understand to send $10 on Ethereum:", 24, GRAY)

items = [
    "What is gas?",
    "What is a gas limit?",
    "What is gas price (gwei)?",
    "What is a base fee?",
    "What is a priority fee (tip)?",
    "What is EIP-1559?",
    "Why did my transaction fail but I still paid the fee?",
    "Why does the same transaction cost $3 now but cost $80 yesterday?",
    "Why do I need ETH to send USDC?",
    "What is a token approval and why does it cost gas?",
    "Why do I need to pay gas to approve, then pay gas again to swap?",
    "What is slippage and why did I lose money?",
    "What is MEV and why did someone front-run my swap?",
]

for i, item in enumerate(items):
    y = 2.4 + i * 0.46
    add_text(slide, 2.2, y, 0.3, 0.4, str(i+1) + ".", 14, RED, True)
    add_text(slide, 2.8, y, 11, 0.4, item, 17, WHITE if i < 6 else GRAY)

add_text(slide, 2, 8.5, 12, 0.5, "On Dina: you send $10. It arrives in 100ms. It costs $0. That's it.", 22, EMERALD, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: Failed transactions still charge
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "THE MOST ABSURD PART", 18, RED, True)
add_text(slide, 2, 2, 12, 2, "On Ethereum, if your\ntransaction fails,\nyou still pay the fee.", 52, WHITE, True)

add_text(slide, 2, 5.5, 12, 1.5, "Imagine going to a restaurant, ordering food, the kitchen\nburns it, and they charge you anyway.\n\nThat's Ethereum's fee model.", 26, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, 2, 7.8, 12, 0.6, 'The industry\'s response: "That\'s by design. The validators still did work."', 20, AMBER)
add_text(slide, 2, 8.3, 12, 0.6, "On Dina: failed transaction = $0 charged. Because fees are $0.", 20, EMERALD, True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: MEV — validators stealing from users
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "MEV: VALIDATORS LITERALLY STEAL FROM USERS", 18, RED, True)

add_text(slide, 2, 1.8, 12, 1, "Maximal Extractable Value — the industry's dirtiest open secret", 24, GRAY)

add_card(slide, 1.5, 2.8, 13, 3.5)
add_text(slide, 1.8, 2.9, 12.5, 0.5, "HOW MEV WORKS", 14, RED, True)
add_multiline(slide, 1.8, 3.4, 12, 3, [
    ("1. You submit a swap: buy Token X for $1,000", WHITE, False, 22),
    ("2. A validator sees your transaction before it's processed", WHITE, False, 22),
    ("3. They buy Token X first (front-running), driving the price up", RED, True, 22),
    ("4. Your transaction executes at the higher price", RED, False, 22),
    ("5. They sell immediately after for profit (sandwich attack)", RED, True, 22),
    ("", WHITE, False, 10),
    ("You paid more. They pocketed the difference.", AMBER, True, 24),
    ("This is legal in crypto. It happens thousands of times per day.", AMBER, False, 20),
])

add_multiline(slide, 2, 6.8, 12, 1.8, [
    ("Total MEV extracted from Ethereum users: $600M+ and counting", RED, True, 22),
    ("", WHITE, False, 8),
    ("These are the same validators the industry says are \"securing\" the network.", GRAY, False, 20),
    ("They're securing the network while picking your pocket.", GRAY, False, 20),
    ("", WHITE, False, 8),
    ("On Dina: 21 known validators run by one company. No MEV. No front-running.", EMERALD, True, 20),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13: The inflation tax
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "THE HIDDEN TAX: TOKEN INFLATION", 18, RED, True)
add_text(slide, 2, 1.7, 12, 1, "Gas fees aren't even enough. They also print new tokens.", 28, WHITE, True)

add_card(slide, 1.5, 3, 13, 2.2)
add_text(slide, 1.8, 3.1, 12.5, 0.5, "ANNUAL TOKEN INFLATION (new tokens printed to pay validators)", 14, RED, True)
add_multiline(slide, 1.8, 3.6, 12, 2, [
    ("Solana:     ~5% per year inflation (decreasing by 15% annually)", WHITE, False, 22),
    ("Ethereum:  ~0.5% per year (offset by EIP-1559 burn — sometimes net deflationary)", WHITE, False, 22),
    ("Avalanche: ~7% per year to staking rewards", WHITE, False, 22),
    ("Polygon:     ~2% per year minted for validators", WHITE, False, 22),
])

add_text(slide, 2, 5.6, 12, 1, "This means: even if you hold tokens and never transact,\nyour purchasing power decreases every year.", 24, AMBER, True)

add_text(slide, 2, 7, 12, 1.5, "So users pay gas fees AND lose value to inflation.\nDouble extraction.", 28, RED, True, PP_ALIGN.CENTER)

add_text(slide, 2, 8.3, 12, 0.6, "Dina has no token. No inflation. Your USDC is your USDC. Always.", 20, EMERALD, True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14: The user experience disaster
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "THE USER EXPERIENCE THEY CREATED", 18, RED, True)
add_text(slide, 2, 1.6, 12, 1, "A normal person trying to use crypto for the first time:", 24, GRAY)

steps = [
    ('"I want to send USDC to my friend"', WHITE),
    ('"I need ETH for gas fees? What\'s ETH?"', WHITE),
    ('"I have to buy ETH on an exchange first?"', WHITE),
    ('"KYC verification? That takes 3 days?"', AMBER),
    ('"OK I bought $20 of ETH for gas"', AMBER),
    ('"Now I need to approve USDC spending? That costs gas too?"', AMBER),
    ('"The approval cost $4. Now I can send."', RED),
    ('"$12 gas fee to send $10??"', RED),
    ('"Transaction failed. I still got charged $12???"', RED),
    ('"I\'m going back to Venmo."', RED),
]

for i, (text, color) in enumerate(steps):
    y = 2.5 + i * 0.55
    num_color = EMERALD if i < 3 else AMBER if i < 6 else RED
    add_text(slide, 2, y, 0.6, 0.5, str(i+1) + ".", 16, num_color, True)
    add_text(slide, 2.7, y, 11, 0.5, text, 19, color)

add_text(slide, 2, 8.3, 12, 0.5, "This is why crypto has 400M wallets but almost no real-world usage.", 22, AMBER, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15: What Visa does
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "THE COMPARISON THEY DON'T WANT YOU TO MAKE", 18, EMERALD, True)

add_card(slide, 1.5, 1.8, 6.2, 3.8)
add_text(slide, 1.8, 1.9, 5.8, 0.5, "VISA", 18, BLUE, True)
add_multiline(slide, 1.8, 2.5, 5.8, 3, [
    ("65,000 TPS", WHITE, True, 24),
    ("~200ms settlement", WHITE, False, 20),
    ("User pays: $0 to swipe", EMERALD, True, 22),
    ("Merchant pays 2-3%", WHITE, False, 18),
    ("", WHITE, False, 10),
    ("Runs on their own servers.", GRAY, False, 18),
    ("Nobody asks who incentivizes", GRAY, False, 18),
    ("Visa's data centers.", GRAY, False, 18),
    ("They just... pay for them.", GRAY, False, 18),
])

add_card(slide, 8.3, 1.8, 6.2, 3.8)
add_text(slide, 8.6, 1.9, 5.8, 0.5, "ETHEREUM", 18, RED, True)
add_multiline(slide, 8.6, 2.5, 5.8, 3, [
    ("~30 TPS", WHITE, True, 24),
    ("~12 minute practical finality", WHITE, False, 20),
    ("User pays: $1-50 per transaction", RED, True, 22),
    ("Plus needs ETH to start", RED, False, 18),
    ("", WHITE, False, 10),
    ("Runs on 1M strangers' servers.", GRAY, False, 18),
    ("Users pay those strangers every", GRAY, False, 18),
    ("single time they transact.", GRAY, False, 18),
    ("And sometimes get front-run.", GRAY, False, 18),
])

add_card(slide, 1.5, 6, 13, 2.2)
add_text(slide, 1.8, 6.1, 12.5, 0.5, "DINA", 18, EMERALD, True)
add_multiline(slide, 1.8, 6.6, 12, 1.5, [
    ("100,000+ TPS capacity  ·  100ms finality  ·  $0.00 per transaction  ·  No token needed", EMERALD, True, 24),
    ("Runs on 21 servers we pay for. Like Visa. But with a public ledger anyone can audit.", WHITE, False, 20),
    ("Visa's speed + blockchain's transparency + nobody pays fees.", WHITE, False, 20),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 16: The ideology trap
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "THE IDEOLOGY TRAP", 18, RED, True)
add_text(slide, 2, 1.7, 12, 1.5, '"Decentralization" became a religion.\nAnd like all religions, questioning it is heresy.', 32, WHITE, True)

add_multiline(slide, 2, 4, 12, 4.5, [
    ("If you say fees are unnecessary:", GRAY, False, 22),
    ('"You don\'t understand decentralization."', RED, False, 24),
    ("", WHITE, False, 10),
    ("If you say tokens are unnecessary:", GRAY, False, 22),
    ('"You don\'t understand crypto economics."', RED, False, 24),
    ("", WHITE, False, 10),
    ("If you say 21 validators is enough:", GRAY, False, 22),
    ('"You don\'t understand security."', RED, False, 24),
    ("", WHITE, False, 10),
    ("Meanwhile:", GRAY, True, 22),
    ("BSC runs on 21 validators and processes billions in daily volume.", EMERALD, False, 22),
    ("Base runs on 1 sequencer and nobody cares.", EMERALD, False, 22),
    ("The ideology only applies when it's convenient.", AMBER, True, 22),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 17: Why nobody fixed it
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "WHY NOBODY FIXED THIS", 18, RED, True)
add_text(slide, 2, 1.8, 12, 1, "Everyone who understands the problem profits from it.", 32, WHITE, True)

add_card(slide, 1.5, 3.2, 13, 5)
add_multiline(slide, 2, 3.4, 12, 4.8, [
    ("Founders won't fix it", RED, True, 26),
    ("Their token holdings are worth billions because of mandatory gas fees.", GRAY, False, 20),
    ("", WHITE, False, 8),
    ("VCs won't fix it", RED, True, 26),
    ("Their returns depend on token price appreciation driven by fee demand.", GRAY, False, 20),
    ("", WHITE, False, 8),
    ("Validators won't fix it", RED, True, 26),
    ("They earn income from gas fees, MEV extraction, and inflation rewards.", GRAY, False, 20),
    ("", WHITE, False, 8),
    ("Exchanges won't fix it", RED, True, 26),
    ("They profit from users buying gas tokens (trading fees on ETH/SOL purchases).", GRAY, False, 20),
    ("", WHITE, False, 8),
    ("Developers won't fix it", AMBER, True, 26),
    ("Most are paid in the native token. Lower token value = lower salary.", GRAY, False, 20),
    ("", WHITE, False, 8),
    ("The entire ecosystem is aligned around keeping fees alive.", AMBER, True, 22),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 18: What it actually costs
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "WHAT A BLOCKCHAIN ACTUALLY COSTS TO RUN", 18, EMERALD, True)
add_text(slide, 2, 1.8, 12, 1, "We know. Because we're running one.", 32, WHITE, True)

add_card(slide, 2, 3, 5, 2.5)
add_text(slide, 2.3, 3.1, 4.5, 0.5, "RIGHT NOW (testnet)", 14, EMERALD, True)
add_text(slide, 2.3, 3.7, 4.5, 0.8, "$75/month", 48, EMERALD, True)
add_text(slide, 2.3, 4.6, 4.5, 0.8, "3 validators on Google Cloud\n100ms blocks, zero fees, real crypto", 16, GRAY)

add_card(slide, 9, 3, 5, 2.5)
add_text(slide, 9.3, 3.1, 4.5, 0.5, "MAINNET (21 validators)", 14, EMERALD, True)
add_text(slide, 9.3, 3.7, 4.5, 0.8, "$504/month", 48, EMERALD, True)
add_text(slide, 9.3, 4.6, 4.5, 0.8, "Full production network\nSame as a mid-tier SaaS subscription", 16, GRAY)

add_card(slide, 2, 6, 12, 2.3)
add_text(slide, 2.3, 6.1, 11.5, 0.5, "FOR CONTEXT", 14, AMBER, True)
add_multiline(slide, 2.3, 6.6, 11, 2, [
    ("Ethereum's network costs: $80B+ in locked staking capital + billions in annual fees", WHITE, False, 20),
    ("Solana's network costs: ~$150M/year in validator hardware + voting fees", WHITE, False, 20),
    ("Dina's network cost: $6,048/year", EMERALD, True, 24),
    ("Same cryptographic verification. Same transaction processing. 99.99% cheaper.", EMERALD, False, 20),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 19: The ask
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 2, 0.8, 12, 1, "WHAT IF WE JUST... STOPPED?", 18, EMERALD, True)

add_text(slide, 2, 2, 12, 2.5, "What if we stopped inventing tokens\nto pay strangers to run servers\nand just ran the servers ourselves?", 44, WHITE, True, PP_ALIGN.CENTER)

add_multiline(slide, 2, 5.5, 12, 3, [
    ("No token. No gas fees. No inflation. No MEV.", EMERALD, True, 28),
    ("No staking. No slashing. No governance theater.", EMERALD, True, 28),
    ("", WHITE, False, 10),
    ("Just USDC. Just transactions. Just 100ms.", WHITE, True, 28),
    ("Just a company paying a $504/month cloud bill.", WHITE, True, 28),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 20: Final
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 2, 1.5, 12, 2, "The blockchain industry built\na toll booth on a free road\nand convinced everyone\nthe tolls were a feature.", 44, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, 2, 5.5, 12, 1, "We removed the toll booth.", 52, WHITE, True, PP_ALIGN.CENTER)

add_text(slide, 2, 7.2, 12, 0.8, "$0.00 per transaction  ·  100ms  ·  4.5% APY  ·  USDC native", 24, EMERALD, True, PP_ALIGN.CENTER)
add_text(slide, 2, 8.2, 12, 0.5, "dina-wallet.web.app                                    Dina Network", 20, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_path = "C:/dina_network/docs/presentations/The_Blockchain_Fee_Machine.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
