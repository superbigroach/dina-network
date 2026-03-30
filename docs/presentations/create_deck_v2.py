"""
The Blockchain Fee Machine v2 — Clean slides + presenter notes
Less text per slide, no overlapping, comprehensive notes
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x34, 0xD3, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LGRAY = RGBColor(0xCB, 0xD5, 0xE1)
RED = RGBColor(0xF8, 0x71, 0x71)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
BLUE = RGBColor(0x60, 0xA5, 0xFA)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def bg(s):
    f = s.background.fill; f.solid(); f.fore_color.rgb = BG

def tx(s, l, t, w, h, txt, sz=24, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.font.name = "Segoe UI"; p.alignment = a
    return tb

def ml(s, l, t, w, h, lines):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, c, b, sz) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.bold = b; p.font.name = "Segoe UI"; p.space_after = Pt(sz * 0.4)
    return tb

def card(s, l, t, w, h, fc=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc; sh.line.fill.background()
    return sh

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

# ═══════════════════════════════════════════════════════════════
# 1: TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 2, 12, 1, "DINA NETWORK", 22, GREEN, True)
tx(s, 2, 3, 12, 2, "The Blockchain\nFee Machine", 58, WHITE, True)
tx(s, 2, 6, 12, 1, "How an industry convinced the world that\nmoving numbers in a database should cost $50.", 26, GRAY)
tx(s, 2, 8, 12, 0.5, "And why they'll never fix it.", 20, RED)
notes(s, """TITLE SLIDE

Open with: "I want to talk about a scam. Not a rug pull. Not a Ponzi scheme. Something much bigger. Something hiding in plain sight in every single blockchain that exists today."

"Every blockchain charges you a fee to move your own money. Ethereum charges $1 to $50. Solana charges fractions of a cent. Even the 'cheap' ones charge something."

"Today I'm going to explain why those fees exist, who profits from them, and why we built a blockchain that charges exactly zero."

Pause. Let it land.

"But first — let's talk about what a transaction actually is."
""")

# ═══════════════════════════════════════════════════════════════
# 2: WHAT A TRANSACTION IS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "WHAT A TRANSACTION ACTUALLY IS", 18, GREEN, True)
tx(s, 2, 2, 12, 1.5, "You want to send $10.", 44, WHITE, True)
ml(s, 2, 4, 10, 3.5, [
    ("Check your signature", WHITE, False, 26),
    ("Check you haven't sent this before", WHITE, False, 26),
    ("Check you have $10", WHITE, False, 26),
    ("Subtract $10 from you", WHITE, False, 26),
    ("Add $10 to them", WHITE, False, 26),
    ("Write it down", WHITE, False, 26),
])
tx(s, 2, 7.8, 12, 0.6, "Total computation: 0.01 milliseconds", 24, GREEN, True)
tx(s, 2, 8.3, 12, 0.6, "Total electricity: $0.000001", 20, GREEN)
notes(s, """WHAT A TRANSACTION IS

"Let me show you what actually happens when you send someone $10 on a blockchain."

"The computer checks your digital signature — is this really from you? Takes about a thousandth of a millisecond."

"Checks the nonce — have we seen this transaction before? Another thousandth of a millisecond."

"Checks your balance — do you actually have $10? Microseconds."

"Subtracts $10 from your account. Adds $10 to theirs. Writes it to disk."

"Total computation time: about one hundredth of a millisecond. The electricity cost is about one millionth of a dollar."

Pause.

"Ethereum charges $1 to $50 for this. That's a markup of one million to fifty million times the actual cost."

"So where does that money go?"
""")

# ═══════════════════════════════════════════════════════════════
# 3: THE MARKUP
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "THE MARKUP", 18, RED, True)
tx(s, 2, 2.5, 12, 1, "Actual cost:", 32, GRAY)
tx(s, 6, 2.5, 6, 1, "$0.000001", 32, GREEN, True)
tx(s, 2, 4, 12, 1, "Ethereum charges:", 32, GRAY)
tx(s, 6, 4, 6, 1, "$1 – $50", 32, RED, True)
tx(s, 2, 6, 12, 1.5, "That's a 50,000,000x markup.", 48, WHITE, True, PP_ALIGN.CENTER)
tx(s, 2, 8, 12, 0.6, "Where does the money go?", 24, AMBER, False, PP_ALIGN.CENTER)
notes(s, """THE MARKUP

"Let that sink in. The actual cost of processing your transaction is one millionth of a dollar."

"Ethereum charges between one and fifty dollars."

"That is a FIFTY MILLION times markup on the actual cost of computation."

"Imagine buying a coffee. The beans cost a tenth of a cent. But the coffee shop charges you fifty thousand dollars. That's the ratio we're talking about."

"No other industry in the world operates at this markup. Not airlines. Not pharmaceuticals. Not even luxury fashion."

"So where does all that money go? It goes to strangers."
""")

# ═══════════════════════════════════════════════════════════════
# 4: PAYING STRANGERS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "WHERE YOUR GAS FEES GO", 18, RED, True)
tx(s, 2, 2.5, 12, 2, "You're paying rent\non strangers' computers.", 48, WHITE, True)
tx(s, 2, 5.5, 12, 2, "That's it. That's the entire reason\ngas fees exist.", 28, GRAY)
notes(s, """PAYING STRANGERS

"Your gas fee doesn't pay for computation. We just established that computation costs basically nothing."

"Your gas fee pays anonymous people around the world to run servers."

"Ethereum has over one million validators — computers run by random people. Those people won't run their computers for free. So every time you move YOUR money, you pay THEM."

"You are paying rent on someone else's computer to update a number in a database. Every. Single. Time."

"Now here's where it gets interesting. You can't just Venmo these strangers. They're anonymous. They're global. So the industry invented an incredibly complex system to pay them."
""")

# ═══════════════════════════════════════════════════════════════
# 5: THE RUBE GOLDBERG MACHINE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 0.6, 12, 0.6, "THE SYSTEM THEY BUILT TO PAY STRANGERS", 18, RED, True)
ml(s, 2, 1.4, 12, 7, [
    ("1.  Invent a token (ETH, SOL, BNB...)", WHITE, False, 24),
    ("2.  Give it a market price", WHITE, False, 24),
    ("3.  Force users to buy it for gas", WHITE, False, 24),
    ("4.  Build a fee market (EIP-1559)", WHITE, False, 24),
    ("5.  Add staking ($80K collateral per validator)", WHITE, False, 24),
    ("6.  Add slashing (punishment for misbehavior)", WHITE, False, 24),
    ("7.  Add inflation (print new tokens when fees aren't enough)", WHITE, False, 24),
    ("8.  Add token burns (control supply from inflation)", WHITE, False, 24),
    ("9.  Add MEV protection (validators front-run users)", WHITE, False, 24),
    ("10. Add governance (vote on fee parameters)", WHITE, False, 24),
    ("11. Add liquid staking (staking is too complex for normal people)", WHITE, False, 24),
    ("12. Add restaking (staking locks up too much capital)", WHITE, False, 24),
])
tx(s, 2, 8.2, 12, 0.6, "12 layers of complexity. All to pay strangers.", 22, AMBER, True, PP_ALIGN.CENTER)
notes(s, """THE RUBE GOLDBERG MACHINE

Go through each one. Slowly. Let the absurdity build.

"Step 1: Invent a token. Every chain has one. ETH, SOL, BNB, AVAX, MATIC."

"Step 2: Get it listed on exchanges so it has a price."

"Step 3: Now force every user to buy this token before they can do anything. That's what gas fees are — mandatory token purchases."

"Step 4: But what if too many people want to transact at once? Build a dynamic fee market. Ethereum's EIP-1559 — base fee goes up and down with demand."

"Step 5: But what if validators cheat? Make them lock up $80,000 as collateral. That's staking."

"Step 6: If they misbehave, take their money. That's slashing."

"Step 7: But fees alone don't cover validator costs. So print new tokens every year. That's inflation. Solana inflates 5% per year."

"Step 8: But inflation devalues the token. So burn some tokens from fees. That's EIP-1559's burn mechanism."

"Step 9: Oh wait, validators discovered they can reorder your transactions to steal from you. That's MEV. Now you need Flashbots and PBS to protect users from their own validators."

"Step 10: Who decides all these parameters? Token holders vote. That's governance."

"Step 11: Staking is too complicated for normal people, so build liquid staking — Lido now controls 28% of all staked ETH. So much for decentralization."

"Step 12: Staking locks up too much capital, so build restaking on top of staking. That's EigenLayer."

"TWELVE layers of complexity. Every single one exists because the first decision was 'let's pay anonymous strangers to run servers.'"
""")

# ═══════════════════════════════════════════════════════════════
# 6: THE REAL REASON FOR TOKENS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "THE REAL REASON TOKENS EXIST", 18, RED, True)
tx(s, 2, 2.5, 12, 1.5, "To make founders rich.", 52, WHITE, True)
card(s, 2, 5, 12, 3)
ml(s, 2.5, 5.2, 11, 2.5, [
    ("Create token → keep 20-40% of supply", WHITE, False, 24),
    ("Require token for gas → creates forced demand", WHITE, False, 24),
    ("Demand pushes price up", WHITE, False, 24),
    ("Founders sell into the demand", RED, True, 28),
])
notes(s, """THE REAL REASON FOR TOKENS

"Now let me tell you the real reason every blockchain has a token. It's not about 'validator incentives' or 'network security.'"

"It's about making the founders obscenely wealthy."

"The playbook is simple. You create a token. You keep 20 to 40 percent of the total supply for founders and early investors."

"Then you launch your blockchain and require everyone to buy your token to pay gas fees. This creates forced, artificial demand."

"That demand pushes the token price up. And the founders sell into that demand."

"This is not speculation. This is the documented tokenomics of every major chain."

"Let me show you the numbers."
""")

# ═══════════════════════════════════════════════════════════════
# 7: FOUNDER WEALTH - SOLANA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "THE NUMBERS", 18, RED, True)
tx(s, 2, 2, 12, 1, "Solana Foundation + Labs", 36, WHITE, True)
tx(s, 2, 3.3, 12, 1, "~300 million SOL", 32, GRAY)
tx(s, 2, 4.5, 12, 1.5, "$15 – $90 billion", 56, RED, True)
tx(s, 2, 6.5, 12, 1, "At peak, worth more than Ford Motor Company.", 24, AMBER)
tx(s, 2, 7.5, 12, 1, "Meanwhile, validators pay $5-10K/month in voting fees\njust to participate in the network.", 20, GRAY)
notes(s, """SOLANA FOUNDER WEALTH

"Solana Foundation and Solana Labs together hold approximately 300 million SOL tokens."

"Depending on market price, that's between 15 and 90 BILLION dollars."

"At its peak, the Solana Foundation's token holdings were worth more than Ford Motor Company. More than FedEx. More than most companies in the Fortune 500."

"And here's the kicker — the validators who actually run the network? They pay $1,000 to $1,800 per month in hardware costs, PLUS $5,000 to $10,000 per month in voting fees just to participate."

"Many Solana validators operate at a loss. They're hoping the SOL token goes up so their staking rewards become worth more."

"The validators are speculating on the token. The founders already have the tokens. The users pay fees that create demand for the token. Everyone is feeding a system designed to enrich the people at the top."
""")

# ═══════════════════════════════════════════════════════════════
# 8: MORE FOUNDERS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "IT'S NOT JUST SOLANA", 18, RED, True)
ml(s, 2, 2, 12, 6, [
    ("Ripple:  ~40B XRP  →  $20-100B", WHITE, True, 32),
    ("Founders among richest people in crypto history", GRAY, False, 18),
    ("", WHITE, False, 14),
    ("BNB:  Binance holds ~80M BNB  →  $30-50B", WHITE, True, 32),
    ("CZ became the richest person in crypto", GRAY, False, 18),
    ("", WHITE, False, 14),
    ("Cardano:  Foundation + IOHK hold ~6B ADA  →  $2-10B", WHITE, True, 32),
    ("Charles Hoskinson: ranch, private jets", GRAY, False, 18),
    ("", WHITE, False, 14),
    ("Avalanche:  Foundation holds ~270M AVAX  →  $3-15B", WHITE, True, 32),
    ("Token sale raised $60M in hours", GRAY, False, 18),
])
notes(s, """MORE FOUNDER WEALTH

Go through each one. Let the numbers speak.

"Ripple — the company holds approximately 40 billion XRP. At various points worth 20 to 100 billion dollars. The co-founder Jed McCaleb alone has sold over a billion dollars worth of XRP."

"BNB — Binance holds roughly 80 million BNB. Worth 30 to 50 billion dollars. CZ became the wealthiest person in crypto entirely from the BNB token that users are forced to buy for gas on BSC."

"Cardano — the Foundation and IOHK together hold about 6 billion ADA. Charles Hoskinson bought a ranch in Colorado, flies private."

"Avalanche — the Foundation holds 270 million AVAX. Their token sale raised 60 million dollars in hours."

"Every single one of these chains charges gas fees. Every single one requires users to buy their native token. And every single founder is a billionaire because of it."

"Why would any of these people eliminate gas fees? It would destroy the value of their holdings."
""")

# ═══════════════════════════════════════════════════════════════
# 9: THE VC GAME
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "THE VC PLAYBOOK", 18, RED, True)
tx(s, 2, 2.5, 12, 1, "Why crypto VCs only fund token projects:", 32, WHITE, True)
card(s, 2, 4, 5.5, 3)
tx(s, 2.3, 4.1, 5, 0.5, "TOKEN PROJECT", 16, RED, True)
ml(s, 2.3, 4.7, 5, 2.2, [
    ("Invest $10M", WHITE, False, 22),
    ("Get tokens at discount", WHITE, False, 22),
    ("Token launches, pumps", WHITE, False, 22),
    ("Sell for $100M", RED, True, 26),
    ("6-12 months", RED, False, 18),
])
card(s, 8.5, 4, 5.5, 3)
tx(s, 8.8, 4.1, 5, 0.5, "REAL COMPANY", 16, GREEN, True)
ml(s, 8.8, 4.7, 5, 2.2, [
    ("Invest $5M for equity", WHITE, False, 22),
    ("Build product, get users", WHITE, False, 22),
    ("Grow revenue over years", WHITE, False, 22),
    ("Exit at $100M", GREEN, True, 26),
    ("5-7 years", GREEN, False, 18),
])
tx(s, 2, 7.8, 12, 0.6, "Same returns. One takes 6 months. Guess which one gets funded.", 22, AMBER, False, PP_ALIGN.CENTER)
notes(s, """THE VC PLAYBOOK

"The founders aren't the only ones profiting. Let me explain the VC game."

"A crypto VC invests $10 million in a token project. They get tokens at a steep discount — maybe $0.10 each. The token launches at $1. They sell for $100 million. 10x return in 6 to 12 months."

"Compare that to investing in a real company. You put in $5 million for equity. The company has to actually build something, get users, grow revenue. Maybe in 5 to 7 years you exit at $100 million."

"Same total return. But one takes 6 months and requires no product-market fit. Just a token launch and some hype."

"Which one do you think VCs prefer?"

"This is why the entire crypto funding ecosystem is designed around tokens. VCs don't care if the technology works. They care if the token pumps."

"And tokens pump when gas fees create mandatory demand."
""")

# ═══════════════════════════════════════════════════════════════
# 10: THE DECENTRALIZATION LIE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, '"DECENTRALIZATION"', 18, RED, True)
tx(s, 2, 2, 12, 1.5, "The word that stops\nall questions.", 48, WHITE, True)
ml(s, 2, 5, 12, 3, [
    ('"Why are fees so high?"', GRAY, False, 24),
    ('"Decentralization costs money."', RED, False, 24),
    ("", WHITE, False, 14),
    ('"Why do I need a token?"', GRAY, False, 24),
    ('"Decentralization requires incentives."', RED, False, 24),
    ("", WHITE, False, 14),
    ('"Why can\'t you just run the servers?"', GRAY, False, 24),
    ('"That would be centralized."', RED, False, 24),
])
notes(s, """THE DECENTRALIZATION LIE

"Now, whenever someone points out the absurdity of gas fees, the industry has a magic word that stops all questions."

Read each exchange slowly. Use a dismissive, condescending tone for the answers — imitating crypto maximalists.

"'Why are fees so high?' — 'Decentralization costs money.'"

"'Why do I need a token to send USDC?' — 'Decentralization requires aligned incentives.'"

"'Why can't you just run the servers yourself?' — 'That would be centralized.'"

"Decentralization has become a religion. And like all religions, questioning it makes you a heretic."

"But let's look at what 'decentralized' actually means for these chains."
""")

# ═══════════════════════════════════════════════════════════════
# 11: THE TRUTH ABOUT DECENTRALIZATION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "HOW 'DECENTRALIZED' ARE THEY REALLY?", 18, RED, True)
ml(s, 2, 2, 12, 6.5, [
    ("Base", BLUE, True, 28),
    ("ONE sequencer. Run by Coinbase. Still charges gas in ETH.", GRAY, False, 20),
    ("", WHITE, False, 12),
    ("BSC", AMBER, True, 28),
    ("21 validators. Binance controls enough BNB to select all of them.", GRAY, False, 20),
    ("", WHITE, False, 12),
    ("Solana", RED, True, 28),
    ("Top 19 validators control the halt threshold. Chain halted 10+ times.", GRAY, False, 20),
    ("Centralized restart every time.", GRAY, False, 20),
    ("", WHITE, False, 12),
    ("Polygon", AMBER, True, 28),
    ("~100 validators. Multisig of known insiders can override the chain.", GRAY, False, 20),
    ("", WHITE, False, 12),
    ("Lido controls 28% of all staked ETH on Ethereum.", RED, True, 20),
    ("One company. 28% of the 'decentralized' network.", RED, False, 20),
])
notes(s, """THE TRUTH ABOUT DECENTRALIZATION

"Let me show you how 'decentralized' these networks actually are."

"Base. Coinbase's L2 chain. Literally one sequencer. One server run by one company. They still charge gas fees in ETH. Why? Because they're an Ethereum L2 and they inherit the fee model. But they could just... not charge fees. They choose to."

"BSC. Binance Smart Chain. 21 validators. Binance controls enough BNB to pick every single validator. It is functionally Binance's private database. With extra steps. And gas fees."

"Solana. The top 19 validators control enough stake to halt the network. And the chain HAS halted — over 10 times. Each time, Solana Labs coordinated a centralized restart. That is the opposite of decentralization."

"Polygon. About 100 validators. But a multisig controlled by known insiders can override the chain."

"And on Ethereum — the most 'decentralized' chain — Lido controls 28 percent of ALL staked ETH. One company. 28 percent of the network."

"'Decentralized' is a marketing word. Not a technical description."
""")

# ═══════════════════════════════════════════════════════════════
# 12: MEV
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "MEV: VALIDATORS STEAL FROM USERS", 18, RED, True)
tx(s, 2, 2.5, 12, 1, "The people 'securing' your transactions\nare also front-running them.", 32, WHITE, True)
ml(s, 2, 5, 12, 3, [
    ("You submit a swap: buy Token X for $1,000", WHITE, False, 24),
    ("Validator sees it, buys Token X first", RED, True, 24),
    ("Your swap executes at a higher price", RED, False, 24),
    ("Validator sells for instant profit", RED, True, 24),
    ("", WHITE, False, 12),
    ("$600M+ extracted from Ethereum users to date.", AMBER, True, 28),
])
notes(s, """MEV — MAXIMAL EXTRACTABLE VALUE

"But wait, it gets worse. Not only do validators charge you fees — they also steal from you."

"It's called MEV — Maximal Extractable Value. Here's how it works."

"You want to swap Token X for $1,000 on a DEX. You submit the transaction."

"A validator sees your transaction BEFORE it gets processed. They see that your buy is going to push the price up. So they buy Token X first, at the lower price."

"Then YOUR transaction executes — at a higher price than it should have been."

"Then the validator immediately sells their tokens at the now-higher price. Instant profit."

"This is called a sandwich attack. It happens thousands of times per day."

"Over $600 million has been extracted from Ethereum users through MEV. These are the same validators that the industry says are 'securing the network.' They're securing it while picking your pocket."

"On Dina, the 21 validators are run by one company with a fiduciary duty to users. No MEV. No front-running. No sandwich attacks."
""")

# ═══════════════════════════════════════════════════════════════
# 13: FAILED TXS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "THE MOST ABSURD PART", 18, RED, True)
tx(s, 2, 2.5, 12, 2.5, "Your transaction fails.\nYou still pay the fee.", 52, WHITE, True)
tx(s, 2, 6, 12, 1.5, "Imagine ordering food at a restaurant.\nThe kitchen burns it.\nThey charge you anyway.", 28, GRAY, False, PP_ALIGN.CENTER)
tx(s, 2, 8.2, 12, 0.6, 'The industry says: "The validators still did work."', 20, AMBER)
notes(s, """FAILED TRANSACTIONS

"Here's my favorite one. On Ethereum, if your transaction FAILS — if it reverts, if it runs out of gas, if anything goes wrong — you STILL pay the gas fee."

"Let me say that again. The transaction didn't work. Your money didn't move. Nothing happened. And you still paid $5, $10, $50."

Pause for effect.

"Imagine going to a restaurant. You order a steak. The kitchen burns it. They throw it in the trash. And they charge you the full price."

"You'd walk out. You'd leave a one-star review. You might call the health department."

"But on Ethereum, this is considered normal. The industry's response is — and I'm quoting — 'The validators still did computational work.'"

"On Dina, a failed transaction costs exactly what a successful transaction costs: zero dollars."
""")

# ═══════════════════════════════════════════════════════════════
# 14: THE USER EXPERIENCE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "WHAT REAL PEOPLE EXPERIENCE", 18, RED, True)
tx(s, 2, 1.8, 12, 0.8, '"I want to send $10 in USDC to my friend"', 28, WHITE)
ml(s, 2, 3, 12, 5.5, [
    ("1.  Wait, I need ETH for gas? What's ETH?", LGRAY, False, 22),
    ("2.  I have to buy ETH on an exchange first?", LGRAY, False, 22),
    ("3.  KYC verification takes 3 days??", AMBER, False, 22),
    ("4.  OK I bought $20 of ETH", AMBER, False, 22),
    ("5.  I need to 'approve' USDC spending? That costs gas too?", AMBER, False, 22),
    ("6.  The approval cost $4", RED, False, 22),
    ("7.  Now $12 gas fee to send $10 of USDC??", RED, False, 22),
    ("8.  Transaction failed. Still charged $12.", RED, True, 22),
    ("9.  I'm going back to Venmo.", RED, True, 26),
])
tx(s, 2, 8.3, 12, 0.5, "400 million crypto wallets. Almost no real-world usage. This is why.", 20, AMBER, True, PP_ALIGN.CENTER)
notes(s, """THE USER EXPERIENCE

Act this out. Roleplay a normal person trying to use crypto for the first time.

"Let me walk you through what a normal person — not a crypto native, just a regular person — experiences when they try to send $10 in USDC to a friend."

Use increasingly frustrated voice:

"'I want to send USDC. Oh, I need ETH for gas fees? What's ETH? I have to buy a DIFFERENT cryptocurrency before I can send the one I have?'"

"'OK, I'll buy ETH. I need to sign up for Coinbase. KYC verification. Upload my driver's license. Wait 3 days.'"

"'Finally, I bought $20 of ETH. Now I can send my USDC. Wait — I need to APPROVE the USDC contract first? That's another transaction? That costs gas too?'"

"'The approval cost $4. Four dollars to give permission to send my own money.'"

"'Now the actual send. $12 gas fee? To send TEN dollars? The fee is MORE than the amount?'"

"'Transaction failed. I STILL GOT CHARGED TWELVE DOLLARS?!'"

"'I'm going back to Venmo.'"

Pause.

"400 million crypto wallets have been created. Almost nobody uses crypto for actual payments. This is why."
""")

# ═══════════════════════════════════════════════════════════════
# 15: WHY NOBODY FIXED IT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "WHY NOBODY HAS FIXED THIS", 18, RED, True)
tx(s, 2, 2, 12, 1.5, "Everyone who understands the problem\nprofits from it.", 40, WHITE, True)
ml(s, 2, 5, 12, 3.5, [
    ("Founders: tokens worth billions because of gas demand", GRAY, False, 22),
    ("VCs: returns depend on token price appreciation", GRAY, False, 22),
    ("Validators: income from fees + MEV + inflation", GRAY, False, 22),
    ("Exchanges: profit from users buying gas tokens", GRAY, False, 22),
    ("Developers: paid in native tokens", GRAY, False, 22),
    ("", WHITE, False, 12),
    ("The entire ecosystem is aligned around keeping fees alive.", AMBER, True, 24),
])
notes(s, """WHY NOBODY FIXED IT

"So why hasn't anyone fixed this? The answer is simple."

"Everyone who understands the problem profits from it."

Go through each one:

"Founders? Their token holdings are worth billions BECAUSE of mandatory gas fees. Gas fees create demand for the token."

"VCs? Their returns depend on the token price. Token price depends on demand. Demand comes from gas fees."

"Validators? They earn income directly from gas fees, MEV extraction, and token inflation."

"Exchanges? Every time someone needs to buy ETH for gas, Coinbase takes a trading fee."

"Developers? Most are paid in the native token. Lower token value equals a lower salary."

"The ENTIRE ecosystem — from top to bottom — is financially aligned around keeping gas fees alive."

"Nobody in the current system has any incentive to eliminate fees. Because fees are the engine that makes everyone at the top wealthy."
""")

# ═══════════════════════════════════════════════════════════════
# 16: WHAT IF WE JUST STOPPED
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 2, 12, 2, "What if we just...\nstopped?", 56, WHITE, True, PP_ALIGN.CENTER)
tx(s, 2, 5.5, 12, 2, "What if we stopped inventing tokens\nand just ran the servers ourselves?", 32, GREEN, False, PP_ALIGN.CENTER)
notes(s, """WHAT IF WE JUST STOPPED

Big pause before this slide. Let the previous slide's message sink in.

"So here's the question I asked myself."

"What if we just... stopped?"

Let it hang.

"What if we stopped inventing tokens to pay strangers? What if we stopped building 12 layers of economic complexity? What if we stopped charging users for computation that costs nothing?"

"What if we just ran the servers ourselves?"

"That's what we built. Let me show you."
""")

# ═══════════════════════════════════════════════════════════════
# 17: INTRODUCING DINA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 0.8, 12, 0.6, "DINA NETWORK", 18, GREEN, True)
tx(s, 2, 2, 12, 1.5, "$0.00 per transaction.\n100ms finality.\nNo token.", 48, WHITE, True)
ml(s, 2, 5.5, 12, 3, [
    ("21 validators on Google Cloud", GREEN, False, 28),
    ("Ed25519 cryptographic verification (same as Solana)", GREEN, False, 28),
    ("Gas metered for safety, priced at $0.00", GREEN, False, 28),
    ("USDC native — no gas token needed", GREEN, False, 28),
])
notes(s, """INTRODUCING DINA

"This is Dina Network."

"Zero dollars per transaction. One hundred millisecond finality. No native token."

"We run 21 validators on Google Cloud. They do the exact same cryptographic verification as Ethereum and Solana — Ed25519 signatures, nonce checking, balance verification, state transitions."

"Gas is metered for safety — smart contracts have limits to prevent infinite loops. But the gas price is zero. Nobody pays for computation."

"The only currency is USDC. No ETH. No SOL. No gas token. You hold USDC, you send USDC, you receive USDC."
""")

# ═══════════════════════════════════════════════════════════════
# 18: THE COST
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "WHAT IT ACTUALLY COSTS", 18, GREEN, True)
card(s, 2, 2.2, 5.5, 3)
tx(s, 2.5, 2.3, 4.5, 0.5, "TESTNET (live now)", 16, GREEN, True)
tx(s, 2.5, 3.2, 4.5, 1, "$75", 64, GREEN, True)
tx(s, 4.5, 3.5, 4.5, 1, "/month", 28, GRAY)
tx(s, 2.5, 4.3, 4.5, 0.8, "3 validators, 100ms blocks\nzero fees, real crypto", 18, GRAY)

card(s, 8.5, 2.2, 5.5, 3)
tx(s, 9, 2.3, 4.5, 0.5, "MAINNET (21 validators)", 16, GREEN, True)
tx(s, 9, 3.2, 4.5, 1, "$504", 64, GREEN, True)
tx(s, 11.5, 3.5, 2, 1, "/month", 28, GRAY)
tx(s, 9, 4.3, 4.5, 0.8, "Full production network\nLess than a SaaS subscription", 18, GRAY)

tx(s, 2, 6, 12, 0.6, "For context:", 20, GRAY)
ml(s, 2, 6.6, 12, 2, [
    ("Ethereum: $80B+ locked capital + billions in annual fees", GRAY, False, 22),
    ("Solana: ~$150M/year in hardware + voting fees", GRAY, False, 22),
    ("Dina: $6,048/year", GREEN, True, 28),
])
notes(s, """WHAT IT COSTS

"Here's what it actually costs to run our network."

"Right now, the testnet is live. Three validators on Google Cloud. $75 per month. Twenty-five dollars per server. 100ms blocks. Zero fees."

"For mainnet with 21 validators: $504 per month. Five hundred and four dollars. That's less than most companies pay for Slack."

Point to the context numbers:

"For context — Ethereum has over $80 billion in locked staking capital plus billions in annual gas fees flowing through the system."

"Solana's validator network costs roughly $150 million per year in hardware and voting fees."

"Dina: six thousand forty-eight dollars per year."

"Same cryptographic verification. Same transaction processing. 99.99% cheaper."
""")

# ═══════════════════════════════════════════════════════════════
# 19: THE 9 WALLETS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 0.8, 12, 0.6, "WHAT WE BUILT", 18, GREEN, True)
tx(s, 2, 1.5, 12, 0.8, "9 wallets per user. All earning yield.", 32, WHITE, True)
card(s, 1.5, 2.8, 4, 2.5)
tx(s, 1.8, 2.9, 3.5, 0.4, "SMART ACCOUNTS (x3)", 14, GREEN, True)
ml(s, 1.8, 3.4, 3.5, 1.8, [
    ("Main  ·  Savings  ·  Backup", WHITE, True, 20),
    ("4.5% APY", GREEN, True, 24),
    ("User controls the keys", GRAY, False, 16),
    ("Each has its own password", GRAY, False, 16),
])

card(s, 6, 2.8, 4, 2.5)
tx(s, 6.3, 2.9, 3.5, 0.4, "AGENT WALLETS (x3)", 14, AMBER, True)
ml(s, 6.3, 3.4, 3.5, 1.8, [
    ("Shopping  ·  Bills  ·  Custom", WHITE, True, 20),
    ("3.5% APY", AMBER, True, 24),
    ("AI signs within daily limits", GRAY, False, 16),
    ("Revocable from any SCA", GRAY, False, 16),
])

card(s, 10.5, 2.8, 4, 2.5)
tx(s, 10.8, 2.9, 3.5, 0.4, "PARALLEL WALLETS (x3)", 14, BLUE, True)
ml(s, 10.8, 3.4, 3.5, 1.8, [
    ("Business  ·  Streaming  ·  API", WHITE, True, 20),
    ("3.0% APY", BLUE, True, 24),
    ("Independent nonces", GRAY, False, 16),
    ("1,000 payments in 100ms", GRAY, False, 16),
])

tx(s, 2, 5.8, 12, 0.5, "All USDC. Users keep 100% of their yield. Dina takes zero.", 22, GREEN, True, PP_ALIGN.CENTER)
ml(s, 2, 6.8, 12, 1.8, [
    ("Real-time yield streaming — balance ticks up every second", WHITE, False, 22),
    ("Balance-based rate limiting — more USDC = more tx capacity", WHITE, False, 22),
    ("$1 minimum balance prevents Sybil attacks", WHITE, False, 22),
    ("Token bucket burst system: up to 50,000 tx burst at $100K+", WHITE, False, 22),
])
notes(s, """THE 9 WALLET SYSTEM

"Let me show you what we actually built."

"Every Dina user gets 9 wallets organized into three tiers."

"Three Smart Contract Accounts — Main, Savings, and Backup. These are your primary wallets. You control the keys. Each has its own password. If one gets compromised, you freeze it from another. They earn 4.5% APY — and you keep ALL of it."

"Three Agent Wallets — for AI. Your shopping bot, your bill payment agent, your custom automation. They operate within daily spending limits. You can revoke them instantly from any Smart Account."

"Three Parallel Wallets — for high throughput. Each has an independent nonce, meaning you can send 1,000 payments simultaneously. Not sequentially. In one 100ms block."

"All earning yield. All in USDC. Users keep 100% of their yield. Dina doesn't take a cut."

"Spam prevention is built into the balance model. Higher balance = more transaction capacity. Token bucket algorithm. $1 minimum to transact. No fees needed."
""")

# ═══════════════════════════════════════════════════════════════
# 20: LIVE DEMO SLIDE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1, 12, 0.6, "IT'S LIVE RIGHT NOW", 18, GREEN, True)
tx(s, 2, 2.5, 12, 1.5, "dina-wallet.web.app", 48, GREEN, True, PP_ALIGN.CENTER)
ml(s, 2, 5, 12, 3, [
    ("Testnet running at 100ms blocks", WHITE, False, 26),
    ("Real Ed25519 signed transactions", WHITE, False, 26),
    ("Real-time yield streaming in your browser", WHITE, False, 26),
    ("Send, receive, convert — all zero fees", WHITE, False, 26),
    ("Tap 'Get 10,000 Test USDC' and try it", WHITE, False, 26),
])
tx(s, 2, 8, 12, 0.5, "Open it. Send a transaction. Watch it confirm in 100ms. Pay $0.", 22, GREEN, True, PP_ALIGN.CENTER)
notes(s, """LIVE DEMO

"And this isn't a pitch deck for something we're going to build. It's live."

"Go to dina-wallet.web.app right now. On your phone, on your laptop."

"Sign in with Google. Tap 'Get 10,000 Test USDC'. Send some to another address. Watch it confirm in 100 milliseconds."

"Zero fees. Real cryptographic signatures. Real blockchain. Running on three servers that cost $75 a month total."

If possible, do a live demo here. Show the wallet. Send a transaction. Show the 100ms confirmation time.

"That's what a blockchain looks like when you remove the toll booth."
""")

# ═══════════════════════════════════════════════════════════════
# 21: CLOSING
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, 2, 1.5, 12, 2, "The blockchain industry\nbuilt a toll booth\non a free road.", 48, GRAY, False, PP_ALIGN.CENTER)
tx(s, 2, 4.5, 12, 2, "And convinced everyone\nthe tolls were a feature.", 48, GRAY, False, PP_ALIGN.CENTER)
tx(s, 2, 7.5, 12, 1, "We removed the toll booth.", 52, WHITE, True, PP_ALIGN.CENTER)
notes(s, """CLOSING

Deliver slowly. Let each line breathe.

"The blockchain industry built a toll booth on a free road."

Pause.

"And convinced everyone the tolls were a feature."

Pause.

"We removed the toll booth."

Long pause.

"$0.00 per transaction. 100 milliseconds. 4.5% APY. USDC native. No token. No gas. No strangers."

"Just a company. Paying a cloud bill. Running a blockchain."

"Thank you."
""")

# Save
out = "C:/dina_network/docs/presentations/The_Blockchain_Fee_Machine_v2.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
